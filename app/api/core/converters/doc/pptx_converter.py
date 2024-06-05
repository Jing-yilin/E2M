from api.core.converters.base_converter import (
    BaseConverter,
)

from api.blueprints.v1.schemas import ResponseData, FileInfo, RequestData

from typing import List
from pathlib import Path
import os

from api.config import Config

# read ppt text
from pptx import Presentation
import os

from dataclasses import dataclass
import json


@dataclass
class SlideElement:
    page: int
    shapes: list
    content: str

    def __repr__(self) -> str:
        return self.__str__()

    def __str__(self) -> str:
        return json.dumps(self.to_dict(), indent=4, ensure_ascii=False)

    def to_dict(self):
        return {"page": self.page, "shapes": self.shapes, "content": self.content}


def get_pptx_slide_elements(pptx_file: str) -> List[SlideElement]:
    prs = Presentation(pptx_file)
    slide_elements = []
    for i, slide in enumerate(prs.slides):
        shapes = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text and shape.text.strip():
                shapes.append(shape.text)
        content = "\n".join(shapes)
        slide_elements.append(SlideElement(i, shapes, content))
    return slide_elements


# logging
import logging

logger = logging.getLogger(__name__)


class PptxConverter(BaseConverter):

    @classmethod
    def allowed_formats(cls) -> list[str]:
        return ["ppt", "pptx"]

    def convert_pptx(
        self,
        **kwargs,
    ) -> ResponseData:

        use_llm = self.request_data.use_llm

        if Config.PPTX_CONVERTER == "default":
            logger.info(f"Converting [{self.file}] with default converter")
            slide_elements = get_pptx_slide_elements(self.file)
            raw = "\n".join([slide.content for slide in slide_elements])

        elif Config.PPTX_CONVERTER == "unstructured":
            logger.info(f"Converting [{self.file}] with unstructured converter")
            from unstructured.partition.pptx import partition_pptx

            pptx_elements = partition_pptx(self.file)
            raw = "\n".join([element.text for element in pptx_elements])
        else:
            raise ValueError(f"Invalid PPTX_CONVERTER: {Config.PPTX_CONVERTER}")

        if Config.ENABLE_LLM and use_llm:
            self.llm_enforce(raw)

        self.set_response_data(status="success", raw=raw)

    def convert_ppt(
        self,
        **kwargs,
    ) -> ResponseData:
        # always use unstructured
        logger.info(f"Converting [{self.file}] with unstructured converter")
        from unstructured.partition.ppt import partition_ppt

        ppt_elements = partition_ppt(self.file)
        raw = "\n".join([element.text for element in ppt_elements])

        if Config.ENABLE_LLM and self.request_data.use_llm:
            self.llm_enforce(raw)

        self.set_response_data(status="success", raw=raw)

    def convert(
        self,
        **kwargs,
    ) -> ResponseData:
        if self.file_info.file_type == "pptx":
            self.convert_pptx(**kwargs)
        elif self.file_info.file_type == "ppt":
            self.convert_ppt(**kwargs)
        else:
            raise ValueError(f"Unsupported file type: {self.file_info.file_type}")

        self.rm_file()
        return self.resp_data
